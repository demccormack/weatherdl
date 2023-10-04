from math import floor
from os import listdir, path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

background_color = RGBColor(0, 34, 102)
text_color = RGBColor(230, 191, 0)


class Briefing:
    """
    Class for creating a weather briefing PPTX based on the supplied configuration.
    """

    def __init__(self, config):
        self._img_dir = config.img_dir
        self._items = config.items

        self._presentation = Presentation()
        self._layout_with_title = self._presentation.slide_layouts[5]
        self._layout_blank = self._presentation.slide_layouts[6]
        for layout in [self._layout_with_title, self._layout_blank]:
            fill = layout.background.fill
            fill.solid()
            fill.fore_color.rgb = background_color

        self._files = listdir(self._img_dir)
        self._current_index = 0
        self._current_slide = None

        self.create()

    def current_img_file(self):
        matching_files = [
            file for file in self._files if file.startswith(f"{self._current_index:03d}")]
        if len(matching_files) > 1:
            raise ValueError(
                f"More than one image with id {self._current_index:03d}: {matching_files}")

        return matching_files[0] if len(
            matching_files) == 1 else None

    def add_slide(self, item, time):
        image_includes_caption = item.get("image_includes_caption")

        slide = self._presentation.slides.add_slide(
            self._layout_blank if image_includes_caption else self._layout_with_title)

        if not image_includes_caption:
            title = slide.shapes.title
            title_text = title.text_frame.paragraphs[0].add_run()
            title_text.font.size = Pt(28)
            title_text.font.bold = True
            title_text.font.color.rgb = text_color
            title.fill.solid()
            title.fill.fore_color.rgb = background_color
            title.line.color.rgb = text_color
            title.line.width = Pt(1.5)
            title_text.text = ' '.join(
                filter(None, [item['name'], time]))

        return slide

    def insert_image(self, file_name):
        pic = self._current_slide.shapes.add_picture(
            path.join(self._img_dir, file_name), Inches(0), Inches(0), width=self._presentation.slide_width)

        # To send the picture to the back, it needs to be repositioned as the first element
        self._current_slide.shapes[0]._element.addprevious(  # pylint: disable=protected-access
            pic._element)  # pylint: disable=protected-access

        aspect_ratio = pic.width / pic.height
        if aspect_ratio < 4 / 3:
            pic.height = self._presentation.slide_height
            pic.width = floor(pic.height * aspect_ratio)
            pic.left = floor(
                (self._presentation.slide_width - pic.width) / 2)
        else:
            max_top_gap = floor(
                self._presentation.slide_height / 4.5)
            space = self._presentation.slide_height - pic.height
            if space > max_top_gap:
                pic.top = max_top_gap

    def set_visibility(self, item, time, file_name):
        show_by_default = item.get("show_by_default") and (item.get(
            "show_by_default") is True or item.get("show_by_default").count(time) > 0)
        if not (file_name and show_by_default):
            self._current_slide._element.set(  # pylint: disable=protected-access
                'show', '0')

    def create(self):
        print("\nBuilding presentation")

        for item in self._items:
            for time in item.get("times", [""]):
                self._current_index += 1

                self._current_slide = self.add_slide(item, time)
                file_name = self.current_img_file()
                self.set_visibility(item, time, file_name)

                if file_name:
                    self.insert_image(file_name)
                else:
                    print(
                        f"No image for {self._current_index} {item['name']} {time}")

    def save_as(self, file_name):
        name = unused_file_name_like(file_name, listdir(self._img_dir))
        full_path = path.join(self._img_dir, name)
        self._presentation.save(full_path)
        print(f"Presentation saved to '{full_path}'")


def unused_file_name_like(name, dir_listing):
    if name not in dir_listing:
        return name

    i = 1
    parts = path.splitext(name)
    while ''.join([f"{parts[0]} ({i})", parts[1]]) in dir_listing:
        i = i + 1

    return ''.join([f"{parts[0]} ({i})", parts[1]])
