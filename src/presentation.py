from math import floor
from os import listdir, path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

background_color = RGBColor(0, 34, 102)
text_color = RGBColor(230, 191, 0)


def create_pptx_from_images(dir, data):
    prs = Presentation()

    layout_with_title = prs.slide_layouts[5]
    layout_blank = prs.slide_layouts[6]
    for layout in [layout_with_title, layout_blank]:
        fill = layout.background.fill
        fill.solid()
        fill.fore_color.rgb = background_color

    index = 0
    files = listdir(dir)
    for item in data:
        for time in item.get("times", [""]):
            index += 1
            matching_files = [
                file for file in files if file.startswith(f"{index:03d}")]
            if len(matching_files) > 1:
                raise Exception(
                    f"More than one image with id {index:03d}: {matching_files}")
            file_name = matching_files[0] if len(
                matching_files) == 1 else False
            image_includes_caption = item.get("image_includes_caption")

            slide = prs.slides.add_slide(
                layout_blank if image_includes_caption else layout_with_title)

            if not image_includes_caption:
                title = slide.shapes.title
                title_text = title.text_frame.paragraphs[0].add_run()
                title_text.font.size = Pt(28)
                title_text.font.bold = True
                title_text.font.color.rgb = text_color
                title.fill.solid()
                title.fill.fore_color.rgb = background_color
                title.line.color.rgb = text_color
                title.line.width = Pt(1.5)
                title_text.text = ' '.join(filter(None, [item['name'], time]))

            if not file_name:
                print(f"No image for {index} {item['name']} {time}")
            else:
                pic = slide.shapes.add_picture(
                    path.join(dir, file_name), Inches(0), Inches(0), width=prs.slide_width)
                # To send the picture to the back, it needs to be repositioned as the first element
                slide.shapes[0]._element.addprevious(pic._element)

                aspect_ratio = pic.width / pic.height
                if aspect_ratio < 4 / 3:
                    pic.height = prs.slide_height
                    pic.width = floor(pic.height * aspect_ratio)
                    pic.left = floor((prs.slide_width - pic.width) / 2)
                else:
                    max_top_gap = floor(prs.slide_height / 4.5)
                    space = prs.slide_height - pic.height
                    if space > max_top_gap:
                        pic.top = max_top_gap

            show_by_default = item.get("show_by_default") and (item.get(
                "show_by_default") == True or item.get("show_by_default").count(time) > 0)
            if not (file_name and show_by_default):
                slide._element.set('show', '0')

    save_safely(prs, path.join(dir, 'briefing.pptx'))

def save_safely(pres, dest):
    dir = path.dirname(dest)
    file = path.basename(dest)
    name = unused_file_name_like(file, listdir(dir))
    pres.save(path.join(dir, name))

def unused_file_name_like(name, dir_listing):
    if name not in dir_listing:
        return name
    
    i = 1
    parts = path.splitext(name)
    while ''.join([f"{parts[0]} ({i})", parts[1]]) in dir_listing:
        i = i + 1
    
    return ''.join([f"{parts[0]} ({i})", parts[1]])
